"""
generate_pptx.py — Creates an editable PowerPoint (.pptx) presentation
for Canva import. 10 slides, GWC brand colors.

Run:  python3 generate_pptx.py
Output: presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── GWC Brand Colors ──────────────────────────────────────────────────────
GREEN       = RGBColor(0x0F, 0x60, 0x3D)   # Primary Green
GOLD        = RGBColor(0xFF, 0xC5, 0x22)   # Primary Gold
DARK_GREEN  = RGBColor(0x03, 0x3F, 0x2B)   # Dark Green
MID_GREEN   = RGBColor(0x1A, 0x99, 0x59)   # Middle Green
BRIGHT      = RGBColor(0xBA, 0xDB, 0x3E)   # Bright Green
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG    = RGBColor(0xF2, 0xF8, 0xF5)
DARK_TEXT   = RGBColor(0x1A, 0x2A, 0x20)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
RED         = RGBColor(0xDC, 0x26, 0x26)
OK_GREEN    = RGBColor(0x16, 0xA3, 0x4A)
AMBER       = RGBColor(0xD9, 0x77, 0x06)

# ── Presentation setup ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)  # 16:9
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]  # blank

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_para(tf, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

def gold_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def footer(slide, num, bg_dark=True):
    tf = add_text(slide, Inches(0.5), H - Inches(0.5), Inches(4), Inches(0.4),
                  "GWC  |  DXHub  |  2026", size=10, color=GRAY if not bg_dark else RGBColor(0x88,0xAA,0x99))
    add_text(slide, W - Inches(1.5), H - Inches(0.5), Inches(1.2), Inches(0.4),
             f"{num} / 10", size=10, color=GRAY if not bg_dark else RGBColor(0x88,0xAA,0x99), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, DARK_GREEN)
gold_bar(s)

add_text(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.8), "GWC", size=52, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(1.7), Inches(6), Inches(0.4), "GOLDEN WEST COLLEGE  |  HUNTINGTON BEACH, CA", size=11, color=RGBColor(0x88,0xBB,0x99))

tf = add_text(s, Inches(0.8), Inches(3.0), Inches(9), Inches(1.5),
              "Automated Student\nEligibility Pipeline", size=38, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(5.0), Inches(9), Inches(0.6),
         "Proactive program-matching that tells students exactly which services\nthey qualify for -- before they even know to ask.", size=15, color=RGBColor(0xCC,0xDD,0xCC))

add_text(s, Inches(0.8), Inches(6.2), Inches(6), Inches(0.4),
         "EOPS  |  CARE  |  CalWORKs  |  and more", size=12, color=GRAY)
footer(s, 1)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()

add_text(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "THE PROBLEM", size=11, bold=True, color=MID_GREEN)
add_text(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "Students don't know what they qualify for", size=30, bold=True, color=DARK_GREEN)

# Stats
for i, (num, lbl) in enumerate([("9,000+","GWC students estimated eligible for EOPS"),("~600","Currently enrolled in EOPS"),("93%","Of eligible students never reach the program")]):
    x = Inches(0.8 + i * 4.0)
    add_text(s, x, Inches(2.2), Inches(3.5), Inches(0.6), num, size=36, bold=True, color=GREEN)
    add_text(s, x, Inches(2.9), Inches(3.5), Inches(0.5), lbl, size=12, color=GRAY)

# Bullets
bullets = [
    "Eligibility criteria scattered across per-program pages",
    "Application abandonment is high -- students start, get confused, stop",
    "Support offices are reactive -- students must find them first",
]
for i, b in enumerate(bullets):
    add_text(s, Inches(1.2), Inches(4.0 + i*0.6), Inches(10), Inches(0.5), f"*  {b}", size=14, color=DARK_TEXT)

footer(s, 2, bg_dark=False)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, DARK_GREEN)
gold_bar(s)

add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), "THE SOLUTION", size=11, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7), "A proactive notification layer -- not an app", size=28, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(1.8), Inches(10), Inches(0.8),
         "When a student logs in, the system matches their Banner + FAFSA + CCCApply data\nagainst every program's criteria automatically.", size=14, color=RGBColor(0xCC,0xDD,0xCC))

# Flow boxes
boxes = ["Banner SIS\nEnrollment | GPA | Residency", "FAFSA / CADAA\nIncome | BOG | Pell | SAI", "CCCApply\nSingle parent | Foster | CalWORKs", "Eligibility Result\nPer program | Instant"]
for i, txt in enumerate(boxes):
    x = Inches(0.8 + i * 3.1)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(2.8), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x7A, 0x4A) if i < 3 else GOLD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE if i < 3 else DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.8),
         '"Congratulations -- you\'re 100% qualified for EOPS and CalWORKs;\nand once we receive one document, you may also qualify for CARE."',
         size=13, color=RGBColor(0x99,0xBB,0xAA))
footer(s, 3)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CONFIRMED vs CONDITIONAL
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()

add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), "HOW IT WORKS", size=11, bold=True, color=MID_GREEN)
add_text(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "Two qualification states -- no ambiguity", size=28, bold=True, color=DARK_GREEN)

# Confirmed card
add_text(s, Inches(1.0), Inches(2.2), Inches(5), Inches(0.4), "100% CONFIRMED", size=14, bold=True, color=OK_GREEN)
add_text(s, Inches(1.0), Inches(2.7), Inches(5), Inches(0.3), "No action needed", size=12, bold=True, color=OK_GREEN)
add_text(s, Inches(1.0), Inches(3.2), Inches(5), Inches(1.2),
         "All criteria verified in Banner, FAFSA, and CCCApply.\nStudent is immediately enrolled or added to priority queue.\n\nExample: EOPS Tier 1 -- income <$19k + homeless + low GPA",
         size=12, color=DARK_TEXT)

# Conditional card
add_text(s, Inches(7.0), Inches(2.2), Inches(5), Inches(0.4), "CONDITIONALLY QUALIFIED", size=14, bold=True, color=AMBER)
add_text(s, Inches(7.0), Inches(2.7), Inches(5), Inches(0.3), "One item needed", size=12, bold=True, color=AMBER)
add_text(s, Inches(7.0), Inches(3.2), Inches(5), Inches(1.2),
         "Core criteria met. One supplemental document required.\nSystem tells student exactly what -- no guessing.\n\nExample: CARE -- needs proof of public assistance letter",
         size=12, color=DARK_TEXT)
footer(s, 4, bg_dark=False)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROGRAMS
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, DARK_GREEN)
gold_bar(s)

add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), "PROGRAMS MATCHED", size=11, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "Three programs live -- built to scale", size=28, bold=True, color=WHITE)

progs = [
    ("EOPS", "Extended Opportunity Programs & Services", "5 criteria: residency, enrollment, units, financial need, ed. qualifier\nTier 1 immediate | Tier 2 waitlist | Tier 3 referral\nCapacity-aware -- never over-enrolls"),
    ("CARE", "Cooperative Agencies Resources for Education", "EOPS-eligible + single parent + dependent child + public assistance\nAdditional grants, childcare referrals, emergency funds"),
    ("CalWORKs", "CA Work Opportunity & Responsibility to Kids", "Enrolled in credit courses + receiving county cash aid\nWork-study, childcare, transportation support, counseling"),
]
for i, (name, full, desc) in enumerate(progs):
    y = Inches(2.0 + i * 1.7)
    add_text(s, Inches(1.0), y, Inches(2), Inches(0.4), name, size=16, bold=True, color=GOLD)
    add_text(s, Inches(1.0), y + Inches(0.4), Inches(4), Inches(0.3), full, size=10, color=RGBColor(0xAA,0xCC,0xBB))
    add_text(s, Inches(5.5), y, Inches(7), Inches(1.2), desc, size=11, color=RGBColor(0xDD,0xEE,0xDD))

add_text(s, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4),
         "Adding a new program = one function. Zero changes to banner, API, or page.", size=11, color=GRAY)
footer(s, 5)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHY NOT AI
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()

add_text(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), "DESIGN DECISION", size=11, bold=True, color=MID_GREEN)
add_text(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "Why rules-based, not AI", size=28, bold=True, color=DARK_GREEN)

# Left column - bad
add_text(s, Inches(1.0), Inches(2.0), Inches(5), Inches(0.4), "Agentic RAG would give us...", size=14, bold=True, color=RED)
for i, txt in enumerate(["Unpredictable -- same student, different result", "Hallucination risk on legal eligibility criteria", "Cost per page load -- every visit charges tokens", "No audit trail -- can't explain decisions"]):
    add_text(s, Inches(1.2), Inches(2.6 + i*0.5), Inches(5), Inches(0.4), f"x  {txt}", size=12, color=DARK_TEXT)

# Right column - good
add_text(s, Inches(7.0), Inches(2.0), Inches(5), Inches(0.4), "Rules engine gives us...", size=14, bold=True, color=OK_GREEN)
for i, txt in enumerate(["Deterministic -- same profile, same result, always", "Free -- zero API cost per student check", "Instant -- runs in milliseconds server-side", "Fully auditable -- every decision has a trace"]):
    add_text(s, Inches(7.2), Inches(2.6 + i*0.5), Inches(5), Inches(0.4), f"*  {txt}", size=12, color=DARK_TEXT)

add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
         "EOPS eligibility is defined by California Title 5. The criteria are fixed and binary.\nAI earns its place for open-ended questions -- not legal checklists.",
         size=11, color=GRAY)
footer(s, 6, bg_dark=False)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DEMO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, GOLD)

add_text(s, Inches(0), Inches(2.0), W, Inches(1.0), "Live Demo", size=44, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(3.2), Inches(9), Inches(0.8),
         "Five student profiles demonstrating confirmed enrollment, waitlist,\nconditional qualification, and referral to alternative programs.",
         size=15, color=DARK_GREEN, align=PP_ALIGN.CENTER)

demos = ["Kylie -- EOPS + CARE + CalWORKs", "Michael -- EOPS Tier 2", "Rosa -- CARE conditional", "Isabella -- Tier 3 referral", "David -- No match"]
for i, d in enumerate(demos):
    x = Inches(1.5 + (i % 3) * 3.6)
    y = Inches(4.5 + (i // 3) * 0.7)
    add_text(s, x, y, Inches(3.3), Inches(0.5), d, size=12, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

footer(s, 7, bg_dark=False)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SCALING
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, DARK_GREEN)
gold_bar(s)

add_text(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "SCALING RESPONSIBLY", size=11, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7), "What happens when 9,000 students qualify?", size=28, bold=True, color=WHITE)

tiers = [("TIER 1", "Immediate", "Highest need enrolled on the spot\nIncome <$19k | Homeless | Foster | GPA <2.0"),
         ("TIER 2", "Waitlist", "Eligible, lower urgency -- managed queue\nEOPS contacts within 2 weeks"),
         ("TIER 3", "Referral", "Program at capacity -- matched to\nCARE, CalWORKs, Basic Needs, Financial Aid")]
for i, (badge, title, desc) in enumerate(tiers):
    x = Inches(0.8 + i * 4.1)
    add_text(s, x, Inches(2.2), Inches(1.5), Inches(0.3), badge, size=10, bold=True, color=GOLD)
    add_text(s, x, Inches(2.6), Inches(3.5), Inches(0.5), title, size=22, bold=True, color=WHITE)
    add_text(s, x, Inches(3.3), Inches(3.5), Inches(1.0), desc, size=11, color=RGBColor(0xCC,0xDD,0xCC))

add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
         "* 9,000 eligible vs 600 served = Chancellor's Office budget justification for increased allocation", size=12, color=RGBColor(0xAA,0xCC,0xBB))
add_text(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.4),
         "* Accept = interest signal into a managed queue -- not instant auto-enrollment", size=12, color=RGBColor(0xAA,0xCC,0xBB))
footer(s, 8)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()

add_text(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "IMPLEMENTATION PATH", size=11, bold=True, color=MID_GREEN)
add_text(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "How GWC plugs this in", size=28, bold=True, color=DARK_GREEN)

add_text(s, Inches(1.0), Inches(2.0), Inches(5), Inches(0.3), "TECHNICAL -- 2 WEEKS", size=11, bold=True, color=GREEN)
tech = ["Read-only access to 6 Banner SIS fields via Ethos API",
        "Read BOG waiver + income from Financial Aid module",
        "Replace mock profile with real auth session (SSO/LMS)",
        "Deploy to GWC subdomain via Vercel"]
for i, t in enumerate(tech):
    add_text(s, Inches(1.2), Inches(2.5 + i*0.45), Inches(5), Inches(0.4), f"*  {t}", size=12, color=DARK_TEXT)

add_text(s, Inches(7.0), Inches(2.0), Inches(5), Inches(0.3), "ORGANIZATIONAL -- 6-10 WEEKS", size=11, bold=True, color=GREEN)
org = ["EOPS director: Banner-verified = auto-enrollment OK",
       "Financial Aid: one data-sharing agreement",
       "IT: one FERPA security review (covers all future students)",
       "Approve the system once -- not 9,000 times"]
for i, t in enumerate(org):
    add_text(s, Inches(7.2), Inches(2.5 + i*0.45), Inches(5.5), Inches(0.4), f"*  {t}", size=12, color=DARK_TEXT)

add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
         "Extension: start EOPS --> add CalWORKs --> add CARE --> add NextUp, DSPS, Golden Promise.\nSame codebase, no rewrite.", size=13, bold=True, color=GREEN)
footer(s, 9, bg_dark=False)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, DARK_GREEN)
gold_bar(s)

add_text(s, Inches(0.8), Inches(1.5), Inches(10), Inches(0.8),
         "Students shouldn't have to find the programs.", size=30, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.4), Inches(10), Inches(0.8),
         "The programs should find the students.", size=30, bold=True, color=GOLD)

stats = [("3", "Programs matched\nEOPS | CARE | CalWORKs"), ("2", "Qualification states\nConfirmed & Conditional"), ("Infinite", "Programs can be added\nOne function each")]
for i, (num, lbl) in enumerate(stats):
    x = Inches(0.8 + i * 4.1)
    add_text(s, x, Inches(4.0), Inches(3), Inches(0.6), num, size=36, bold=True, color=GOLD)
    add_text(s, x, Inches(4.7), Inches(3.5), Inches(0.6), lbl, size=11, color=RGBColor(0xCC,0xDD,0xCC))

add_text(s, Inches(0.8), Inches(6.2), Inches(10), Inches(0.4),
         "github.com/miaaoyama/GoldenWestBanner  |  Next.js | TypeScript | DXHub design", size=11, color=GRAY)
footer(s, 10)

# ── Save ──────────────────────────────────────────────────────────────────
OUTPUT = "/Users/marinapmoreira/Downloads/GoldenWestBanner/presentation.pptx"
prs.save(OUTPUT)
print(f"Done! Saved to: {OUTPUT}")
